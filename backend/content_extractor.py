"""Content extraction from various sources: YouTube, URLs, PDFs, DOCX, PPT, EPUB, images (OCR), audio (Whisper), text."""
import io
import re
import os
import tempfile
from typing import Optional
from urllib.parse import urlparse, parse_qs

import httpx
from bs4 import BeautifulSoup
from pypdf import PdfReader
import docx as docx_lib
from pptx import Presentation
from ebooklib import epub
import ebooklib
from PIL import Image
import pytesseract
import yt_dlp

def _yt_id(url: str) -> Optional[str]:
    p = urlparse(url)
    if p.hostname in ("youtu.be",):
        return p.path.lstrip("/")
    if p.hostname and "youtube.com" in p.hostname:
        if p.path == "/watch":
            return parse_qs(p.query).get("v", [None])[0]
        m = re.match(r"^/(embed|shorts)/([^/?]+)", p.path)
        if m:
            return m.group(2)
    return None

def extract_from_youtube(url: str) -> dict:
    vid = _yt_id(url)
    if not vid:
        raise ValueError("URL YouTube non valido")
    try:
        ydl_opts = {
            'skip_download': True,
            'writesubtitles': True,
            'writeautomaticsub': True,
            'subtitleslangs': ['it', 'en'],
            'quiet': True,
            'no_warnings': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            subs = info.get('subtitles') or info.get('automatic_captions')
            if not subs:
                raise ValueError("Nessun sottotitolo disponibile per questo video")
            
            lang = 'it' if 'it' in subs else ('en' if 'en' in subs else list(subs.keys())[0])
            track = next((t for t in subs[lang] if t['ext'] == 'json3'), subs[lang][0])
            
            with httpx.Client(follow_redirects=True, timeout=10.0) as c:
                r = c.get(track['url'])
                r.raise_for_status()
                if 'json' in track['ext']:
                    data = r.json()
                    text = "".join([seg.get('utf8', '') for event in data.get('events', []) for seg in event.get('segs', [])])
                    text = " ".join(text.split())
                else:
                    text = r.text
    except Exception as e:
        raise ValueError(f"Estrazione bloccata da YouTube: {e}")
    return {"text": text, "source_type": "youtube", "title": f"YouTube: {vid}"}


def extract_from_url(url: str) -> dict:
    if "youtube.com" in url or "youtu.be" in url:
        return extract_from_youtube(url)
    with httpx.Client(follow_redirects=True, timeout=30.0, headers={"User-Agent": "Mozilla/5.0"}) as c:
        r = c.get(url)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        title = (soup.title.string.strip() if soup.title and soup.title.string else url)
        text = " ".join(soup.get_text(separator=" ").split())
    return {"text": text, "source_type": "url", "title": title}


def extract_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((p.extract_text() or "") for p in reader.pages)


def extract_from_docx(data: bytes) -> str:
    doc = docx_lib.Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        chunks.append("\n".join(parts))
    return "\n\n".join(chunks)


def extract_from_epub(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as f:
        f.write(data)
        path = f.name
    try:
        book = epub.read_epub(path)
        chunks = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "lxml")
                chunks.append(soup.get_text(separator=" "))
        return "\n\n".join(chunks)
    finally:
        os.unlink(path)


def extract_from_image(data: bytes) -> str:
    raise ValueError("L'estrazione del testo dalle immagini non è supportata in questa versione (hosting su Vercel).")


def extract_from_audio(data: bytes, filename: str, api_key: str) -> str:
    """Audio transcription non disponibile in questa versione."""
    raise ValueError("La trascrizione audio non è supportata. Carica un file PDF, DOCX, PPTX, EPUB o testo.")



def extract_from_file(filename: str, data: bytes, api_key: str) -> dict:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return {"text": extract_from_pdf(data), "source_type": "pdf", "title": filename}
    if ext in (".docx", ".doc"):
        return {"text": extract_from_docx(data), "source_type": "docx", "title": filename}
    if ext in (".pptx", ".ppt"):
        return {"text": extract_from_pptx(data), "source_type": "pptx", "title": filename}
    if ext == ".epub":
        return {"text": extract_from_epub(data), "source_type": "epub", "title": filename}
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"):
        return {"text": extract_from_image(data), "source_type": "image", "title": filename}
    if ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".mp4", ".mpeg", ".webm"):
        return {"text": extract_from_audio(data, filename, api_key), "source_type": "audio", "title": filename}
    if ext in (".txt", ".md"):
        return {"text": data.decode("utf-8", errors="ignore"), "source_type": "text", "title": filename}
    raise ValueError(f"Formato non supportato: {ext}")
